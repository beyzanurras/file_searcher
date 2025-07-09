import sys
import os
from PyQt5.QtWidgets import (
    QApplication, QMainWindow, QWidget, QVBoxLayout, QHBoxLayout,
    QLabel, QLineEdit, QPushButton, QListWidget, QFileDialog, QStatusBar, QCheckBox, QGroupBox, QMenu, QTextEdit, QGridLayout, QSplitter, QListWidgetItem, QMessageBox, QRadioButton, QButtonGroup
)
from PyQt5.QtCore import QThread, pyqtSignal, Qt, QEvent, QSettings
from PyQt5.QtGui import QCursor
from PyQt5.QtGui import QFont, QCursor, QTextCharFormat, QTextCursor, QColor
from file_searcher import FileSearcher
import subprocess
import platform
import zipfile
import xml.etree.ElementTree as ET
import io
import multiprocessing
from multiprocessing import Pool, Manager
import fitz
import re

# Yardımcı: Dosya türü kategorileri
TXT_EXTS = ['.txt']
PDF_EXTS = ['.pdf']
OFFICE_EXTS = ['.docx', '.docm', '.dotx', '.dotm', '.xlsx', '.xlsm', '.xltx', '.xltm', '.xlsb', '.pptx', '.pptm', '.ppsx', '.ppsm', '.potx', '.potm', '.vsdx', '.vsd']

def matches_keyword_simple(text, keyword, match_type, case_sensitive):
    """UI'dan bağımsız kelime eşleştirme fonksiyonu"""
    import re
    
    # Büyük/küçük harf ayarı
    search_text = text if case_sensitive else text.lower()
    search_keyword = keyword if case_sensitive else keyword.lower()
    
    if match_type == 0:  # Normal arama
        return search_keyword in search_text
    elif match_type == 1:  # Tam kelime
        # Word boundary kullanarak tam kelime ara
        flags = 0 if case_sensitive else re.IGNORECASE
        pattern = r'\b' + re.escape(keyword) + r'\b'
        return bool(re.search(pattern, text, flags))
    elif match_type == 2:  # Başlangıç
        # Kelime başında ara
        flags = 0 if case_sensitive else re.IGNORECASE
        pattern = r'\b' + re.escape(keyword)
        return bool(re.search(pattern, text, flags))
    elif match_type == 3:  # Bitiş
        # Kelime sonunda ara
        flags = 0 if case_sensitive else re.IGNORECASE
        pattern = re.escape(keyword) + r'\b'
        return bool(re.search(pattern, text, flags))
    
    return False

def file_search_worker(args):
    file_path, keyword_list, extensions, case_sensitive, match_type = args
    file_extension = os.path.splitext(file_path)[1].lower()
    try:
        content = ""
        if file_extension == '.txt':
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
        elif file_extension in ['.docx', '.docm', '.dotx', '.dotm']:
            from docx import Document
            doc = Document(file_path)
            content = '\n'.join([p.text for p in doc.paragraphs])
        elif file_extension in ['.xlsx', '.xlsm', '.xltx', '.xltm']:
            from openpyxl import load_workbook
            wb = load_workbook(file_path, data_only=True)
            texts = []
            for sheet in wb.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    for cell in row:
                        if cell is not None:
                            texts.append(str(cell))
            content = '\n'.join(texts)
        elif file_extension == '.xlsb':
            try:
                from pyxlsb import open_workbook
                texts = []
                with open_workbook(file_path) as wb:
                    for sheet_name in wb.get_sheet_names():
                        with wb.get_sheet(sheet_name) as sheet:
                            for row in sheet.rows():
                                for cell in row:
                                    if cell.v is not None:
                                        texts.append(str(cell.v))
                content = '\n'.join(texts)
            except ImportError:
                return None
        elif file_extension in ['.pptx', '.pptm', '.ppsx', '.ppsm', '.potx', '.potm']:
            try:
                from pptx import Presentation
                prs = Presentation(file_path)
                texts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text_frame") and shape.text_frame:
                            for paragraph in shape.text_frame.paragraphs:
                                if paragraph.text.strip():
                                    texts.append(paragraph.text.strip())
                content = '\n'.join(texts)
            except ImportError:
                return None
        elif file_extension == '.pdf':
            with fitz.open(file_path) as doc:
                content = '\n'.join(page.get_text() for page in doc)
        elif file_extension == '.vsdx':
            metinler = read_vsdx_text(file_path)
            content = ' '.join(metinler)
        # .vsd için içerik okuma yok
        for keyword in keyword_list:
            if matches_keyword_simple(content, keyword, match_type, case_sensitive):
                return file_path
    except Exception:
        return None
    return None

class SearchThread(QThread):
    dosya_bulundu = pyqtSignal(str)
    arama_bitti = pyqtSignal(int)
    arama_durumu = pyqtSignal(str)

    def __init__(self, directory, keywords, extensions, case_sensitive=False, match_type=0):
        super().__init__()
        self.directory = directory
        self.keywords = keywords
        self.extensions = extensions
        self.case_sensitive = case_sensitive
        self.match_type = match_type
        self._stop_requested = False

    def run(self):
        self.arama_durumu.emit("Arama yapılıyor...")
        keyword_list = [k.strip() for k in self.keywords.split(',') if k.strip()]
        if not keyword_list:
            self.arama_durumu.emit("Lütfen aranacak kelimeleri girin.")
            self.arama_bitti.emit(0)
            return
        if not self.extensions:
            self.arama_durumu.emit("Lütfen en az bir dosya türü seçin.")
            self.arama_bitti.emit(0)
            return
        # 1. Tüm dosya yollarını topla ve kategorilere ayır
        txt_files, office_files, pdf_files = [], [], []
        for root, dirs, files in os.walk(self.directory):
            if self._stop_requested:
                self.arama_durumu.emit("Arama iptal edildi.")
                self.arama_bitti.emit(0)
                return
            for file in files:
                file_path = os.path.join(root, file)
                ext = os.path.splitext(file_path)[1].lower()
                if ext not in self.extensions:
                    continue
                if ext in TXT_EXTS:
                    txt_files.append(file_path)
                elif ext in OFFICE_EXTS:
                    office_files.append(file_path)
                elif ext in PDF_EXTS:
                    pdf_files.append(file_path)
        # 2. Paralel arama fonksiyonu
        def parallel_search(file_list):
            found = []
            if not file_list:
                return found
            with Pool(processes=max(1, multiprocessing.cpu_count()-1)) as pool:
                args = [(fp, keyword_list, self.extensions, self.case_sensitive, self.match_type) for fp in file_list]
                for result in pool.imap_unordered(file_search_worker, args):
                    if self._stop_requested:
                        pool.terminate()
                        break
                    if result:
                        self.dosya_bulundu.emit(result)
                        found.append(result)
            return found
        toplam_bulunan = 0
        for file_list in [txt_files, office_files, pdf_files]:
            if self._stop_requested:
                break
            bulunanlar = parallel_search(file_list)
            toplam_bulunan += len(bulunanlar)
        if self._stop_requested:
            self.arama_durumu.emit("Arama iptal edildi.")
        else:
            self.arama_durumu.emit(f"Arama tamamlandı. {toplam_bulunan} dosya bulundu.")
        self.arama_bitti.emit(toplam_bulunan)

    def stop(self):
        self._stop_requested = True

class MainWindow(QMainWindow):
    def __init__(self):
        super().__init__()
        self.setWindowTitle("Dosya İçeriği Arama Motoru")
        self.setGeometry(200, 200, 750, 700)
        self.search_thread = None
        self.settings = QSettings("Beyza", "DosyaAramaUygulamasi")
        self.init_ui()

    def init_ui(self):
        # Ana widget ve ana layout
        central = QWidget()
        self.setCentralWidget(central)
        main_layout = QVBoxLayout(central)
        main_layout.setContentsMargins(30, 30, 30, 20)
        main_layout.setSpacing(18)

        # --- MODERN QSS TEMA ---
        self.setStyleSheet('''
        QWidget {
            background: #f6f7fb;
            font-family: 'Segoe UI', 'Arial', sans-serif;
            font-size: 15px;
        }
        QLabel#TitleLabel {
            font-size: 26px;
            font-weight: 700;
            color: #222;
            letter-spacing: 1px;
        }
        QLabel {
            color: #444;
        }
        QLineEdit {
            background: #fff;
            border: 1.5px solid #d0d4e4;
            border-radius: 8px;
            padding: 7px 12px;
            font-size: 15px;
        }
        QLineEdit:focus {
            border: 1.5px solid #1976d2;
            background: #f0f6ff;
        }
        QPushButton {
            background: qlineargradient(x1:0, y1:0, x2:1, y2:0, stop:0 #1976d2, stop:1 #42a5f5);
            color: white;
            font-weight: 600;
            font-size: 16px;
            border: none;
            border-radius: 8px;
            padding: 9px 22px;
            margin: 0 4px;
        }
        QPushButton:hover {
            background: qlineargradient(x1:0, y1:0, x2:1, y2:0, stop:0 #1565c0, stop:1 #64b5f6);
        }
        QPushButton:pressed {
            background: #0d47a1;
        }
        QCheckBox {
            spacing: 8px;
            font-size: 14px;
        }
        QCheckBox::indicator {
            width: 20px;
            height: 20px;
            border-radius: 6px;
            border: 1.5px solid #b0b8c9;
            background: #fff;
        }
        QCheckBox::indicator:checked {
            background: #1976d2;
            border: 1.5px solid #1976d2;
        }
        QCheckBox::indicator:hover {
            border: 1.5px solid #1976d2;
        }
        QListWidget {
            background: #fff;
            border: 1.5px solid #d0d4e4;
            border-radius: 10px;
            font-size: 15px;
            padding: 6px;
        }
        QListWidget::item:selected {
            background: #e3f0fd;
            color: #1976d2;
            border-radius: 6px;
        }
        QTextEdit {
            background: #fff;
            border: 1.5px solid #d0d4e4;
            border-radius: 10px;
            font-size: 15px;
            padding: 10px;
        }
        QStatusBar {
            background: #f0f2f7;
            color: #1976d2;
            border-top: 1px solid #d0d4e4;
            font-size: 14px;
        }
        ''')

        # --- BAŞLIK ---
        title = QLabel("Dosya İçeriği Arama Motoru")
        title.setObjectName("TitleLabel")
        title.setAlignment(Qt.AlignCenter)
        main_layout.addWidget(title)

        # --- Dizin seçimi ---
        dir_layout = QHBoxLayout()
        dir_label = QLabel("Aranacak Dizin:")
        dir_label.setMinimumWidth(120)
        self.dir_edit = QLineEdit()
        self.dir_edit.setReadOnly(True)
        self.dir_edit.setMinimumHeight(36)
        dir_btn = QPushButton("Dizin Seç")
        dir_btn.setMinimumHeight(36)
        dir_btn.clicked.connect(self.select_directory)
        dir_layout.addWidget(dir_label)
        dir_layout.addWidget(self.dir_edit)
        dir_layout.addWidget(dir_btn)
        main_layout.addLayout(dir_layout)

        # --- Arama metni ---
        word_layout = QHBoxLayout()
        word_layout.addWidget(QLabel("Aranacak Kelimeler (virgülle ayırın):"))
        self.word_edit = QLineEdit()
        self.word_edit.setMinimumHeight(35)
        self.word_edit.setStyleSheet("font-size: 15px; padding: 8px; border: 2px solid #ddd; border-radius: 8px; background: white;")
        word_layout.addWidget(self.word_edit)
        main_layout.addLayout(word_layout)
        
        # --- Arama seçenekleri ---
        options_layout = QHBoxLayout()
        
        # Büyük/küçük harf duyarlılığı
        self.case_sensitive_cb = QCheckBox("Büyük/küçük harf duyarlı")
        self.case_sensitive_cb.setStyleSheet("font-size: 14px; margin: 5px;")
        options_layout.addWidget(self.case_sensitive_cb)
        
        # Kelime eşleştirme seçenekleri
        match_group = QGroupBox("Kelime Eşleştirme:")
        match_group.setStyleSheet("QGroupBox { font-weight: bold; font-size: 14px; margin: 5px; } QGroupBox::title { color: #1976d2; }")
        match_layout = QHBoxLayout(match_group)
        match_layout.setSpacing(15)  # Radio buttonlar arası boşluk
        
        # Radio button grubu oluştur
        self.match_button_group = QButtonGroup()
        
        self.normal_match_rb = QRadioButton("Normal")
        self.exact_match_rb = QRadioButton("Tam kelime")
        self.starts_with_rb = QRadioButton("Başlangıç")
        self.ends_with_rb = QRadioButton("Bitiş")
        
        # Radio button stilini ayarla
        radio_style = "QRadioButton { font-size: 13px; margin: 2px; min-width: 80px; }"
        self.normal_match_rb.setStyleSheet(radio_style)
        self.exact_match_rb.setStyleSheet(radio_style)
        self.starts_with_rb.setStyleSheet(radio_style)
        self.ends_with_rb.setStyleSheet(radio_style)
        
        # Varsayılan olarak normal seçili
        self.normal_match_rb.setChecked(True)
        
        # Grup içine ekle
        self.match_button_group.addButton(self.normal_match_rb, 0)
        self.match_button_group.addButton(self.exact_match_rb, 1)
        self.match_button_group.addButton(self.starts_with_rb, 2)
        self.match_button_group.addButton(self.ends_with_rb, 3)
        
        # Layout'a ekle
        match_layout.addWidget(self.normal_match_rb)
        match_layout.addWidget(self.exact_match_rb)
        match_layout.addWidget(self.starts_with_rb)
        match_layout.addWidget(self.ends_with_rb)
        
        options_layout.addWidget(match_group)
        options_layout.addStretch()
        main_layout.addLayout(options_layout)

        # --- Dosya türü seçim kutuları ---
        filetype_layout = QVBoxLayout()
        filetype_layout.setSpacing(6)
        filetype_layout.setContentsMargins(0, 0, 0, 0)
        def section_row(title, checkboxes):
            row = QHBoxLayout()
            row.setSpacing(10)
            label = QLabel(title)
            label.setStyleSheet("font-weight: 600; font-size: 14px; color: #1976d2; margin-right: 8px;")
            row.addWidget(label)
            for cb in checkboxes:
                cb.setStyleSheet("font-size: 14px; margin-right: 2px; margin-left: 2px;")
                cb.setChecked(True)
                row.addWidget(cb)
            row.addStretch(1)
            return row
        self.cb_docx = QCheckBox(".docx")
        self.cb_docm = QCheckBox(".docm")
        self.cb_dotx = QCheckBox(".dotx")
        self.cb_dotm = QCheckBox(".dotm")
        filetype_layout.addLayout(section_row("Word:", [self.cb_docx, self.cb_docm, self.cb_dotx, self.cb_dotm]))
        self.cb_xlsx = QCheckBox(".xlsx")
        self.cb_xlsm = QCheckBox(".xlsm")
        self.cb_xltx = QCheckBox(".xltx")
        self.cb_xltm = QCheckBox(".xltm")
        self.cb_xlsb = QCheckBox(".xlsb")
        filetype_layout.addLayout(section_row("Excel:", [self.cb_xlsx, self.cb_xlsm, self.cb_xltx, self.cb_xltm, self.cb_xlsb]))
        self.cb_pptx = QCheckBox(".pptx")
        self.cb_pptm = QCheckBox(".pptm")
        self.cb_ppsx = QCheckBox(".ppsx")
        self.cb_ppsm = QCheckBox(".ppsm")
        self.cb_potx = QCheckBox(".potx")
        self.cb_potm = QCheckBox(".potm")
        filetype_layout.addLayout(section_row("PowerPoint:", [self.cb_pptx, self.cb_pptm, self.cb_ppsx, self.cb_ppsm, self.cb_potx, self.cb_potm]))
        self.cb_txt = QCheckBox(".txt")
        self.cb_pdf = QCheckBox(".pdf")
        filetype_layout.addLayout(section_row("Diğer:", [self.cb_txt, self.cb_pdf]))
        # Visio
        self.cb_vsdx = QCheckBox(".vsdx")
        self.cb_vsd = QCheckBox(".vsd")
        filetype_layout.addLayout(section_row("Visio:", [self.cb_vsdx, self.cb_vsd]))
        main_layout.addLayout(filetype_layout)

        # --- Butonlar ---
        btn_layout = QHBoxLayout()
        self.search_btn = QPushButton("Aramayı Başlat")
        self.search_btn.setMinimumHeight(40)
        self.search_btn.clicked.connect(self.toggle_search)
        btn_layout.addWidget(self.search_btn)
        self.save_btn = QPushButton("Sonuçları Kaydet")
        self.save_btn.setMinimumHeight(40)
        self.save_btn.clicked.connect(self.save_results)
        btn_layout.addWidget(self.save_btn)
        main_layout.addLayout(btn_layout)

        # --- Sonuç listesi ---
        self.result_list = QListWidget()
        self.result_list.setMinimumHeight(180)
        main_layout.addWidget(self.result_list)
        self.result_list.itemDoubleClicked.connect(self.open_selected_file)
        self.result_list.setContextMenuPolicy(Qt.CustomContextMenu)
        self.result_list.customContextMenuRequested.connect(self.show_context_menu)
        # Eski önizleme bağlantıları kaldırıldı

        # --- Alt bölüm: İki sütunlu yapı ---
        alt_widget = QWidget()
        alt_layout = QHBoxLayout(alt_widget)
        
        # Sol taraf: Dosya listesi
        sol_widget = QWidget()
        sol_layout = QVBoxLayout(sol_widget)
        sol_layout.addWidget(QLabel("📁 Bulunan Dosyalar:"))
        self.dosya_listesi = QListWidget()
        self.dosya_listesi.setMinimumHeight(150)
        sol_layout.addWidget(self.dosya_listesi)
        
        # Sağ taraf: Satır bilgileri
        sag_widget = QWidget()
        sag_layout = QVBoxLayout(sag_widget)
        sag_layout.addWidget(QLabel("📍 Satır Bilgileri:"))
        self.satir_bilgileri = QListWidget()
        self.satir_bilgileri.setMinimumHeight(150)
        sag_layout.addWidget(self.satir_bilgileri)
        
        # QSplitter ile iki tarafı ayır
        splitter = QSplitter(Qt.Horizontal)
        splitter.addWidget(sol_widget)
        splitter.addWidget(sag_widget)
        splitter.setStretchFactor(0, 1)
        splitter.setStretchFactor(1, 1)
        
        alt_layout.addWidget(splitter)
        main_layout.addWidget(alt_widget)
        
        # Dosya listesi seçimi ile satır bilgilerini senkronize et
        self.dosya_listesi.currentItemChanged.connect(self.dosya_secildi)

        # --- Durum çubuğu ---
        self.status_bar = QStatusBar()
        self.setStatusBar(self.status_bar)
        self.status_bar.showMessage("Lütfen bir dizin seçin.")

        # Son kaydedilen dizini yükle
        last_dir = self.settings.value("last_directory", "")
        if last_dir:
            self.dir_edit.setText(last_dir)

        # Genel pencere arka planı
        self.setStyleSheet("QWidget { background: #f4f6fa; } QLabel { font-size: 15px; }")

        self._searching = False

    def select_directory(self):
        folder = QFileDialog.getExistingDirectory(self, "Dizin Seç")
        if folder:
            self.dir_edit.setText(folder)
            self.status_bar.showMessage("Dizin seçildi. Aranacak kelimeleri girin.")
            # Ayarlara kaydet
            self.settings.setValue("last_directory", folder)

    def toggle_search(self):
        if not self._searching:
            self.start_search()
        else:
            self.stop_search()

    def start_search(self):
        directory = self.dir_edit.text().strip()
        keywords = self.word_edit.text().strip()
        extensions = []
        
        # Word dosyaları
        if self.cb_docx.isChecked():
            extensions.append('.docx')
        if self.cb_docm.isChecked():
            extensions.append('.docm')
        if self.cb_dotx.isChecked():
            extensions.append('.dotx')
        if self.cb_dotm.isChecked():
            extensions.append('.dotm')
            
        # Excel dosyaları
        if self.cb_xlsx.isChecked():
            extensions.append('.xlsx')
        if self.cb_xlsm.isChecked():
            extensions.append('.xlsm')
        if self.cb_xltx.isChecked():
            extensions.append('.xltx')
        if self.cb_xltm.isChecked():
            extensions.append('.xltm')
        if self.cb_xlsb.isChecked():
            extensions.append('.xlsb')
            
        # PowerPoint dosyaları
        if self.cb_pptx.isChecked():
            extensions.append('.pptx')
        if self.cb_pptm.isChecked():
            extensions.append('.pptm')
        if self.cb_ppsx.isChecked():
            extensions.append('.ppsx')
        if self.cb_ppsm.isChecked():
            extensions.append('.ppsm')
        if self.cb_potx.isChecked():
            extensions.append('.potx')
        if self.cb_potm.isChecked():
            extensions.append('.potm')
            
        # Diğer dosya türleri
        if self.cb_txt.isChecked():
            extensions.append('.txt')
        if self.cb_pdf.isChecked():
            extensions.append('.pdf')
            
        # Visio dosyaları
        if self.cb_vsdx.isChecked():
            extensions.append('.vsdx')
        if self.cb_vsd.isChecked():
            extensions.append('.vsd')

        if not directory:
            self.status_bar.showMessage("Lütfen bir dizin seçin.")
            return
        if not keywords:
            self.status_bar.showMessage("Lütfen aranacak kelimeleri girin.")
            return
        if not extensions:
            self.status_bar.showMessage("Lütfen en az bir dosya türü seçin.")
            return
        self.result_list.clear()
        self.dosya_listesi.clear() # Dosya listesini temizle
        self.satir_bilgileri.clear() # Satır bilgilerini temizle
        self.status_bar.showMessage("Arama yapılıyor...")
        self.search_btn.setText("Aramayı Durdur")
        self.search_btn.setStyleSheet("background-color: #d32f2f; color: white; font-weight: bold; font-size: 16px; border-radius: 8px;")
        self.search_btn.setEnabled(True)
        self._searching = True
        case_sensitive = self.case_sensitive_cb.isChecked()
        match_type = self.match_button_group.checkedId()
        self.search_thread = SearchThread(directory, keywords, extensions, case_sensitive, match_type)
        self.search_thread.dosya_bulundu.connect(self.add_result)
        self.search_thread.arama_bitti.connect(self.search_finished)
        self.search_thread.arama_durumu.connect(self.status_bar.showMessage)
        self.search_thread.start()

    def stop_search(self):
        if self.search_thread and self._searching:
            self.search_thread.stop()
            self.status_bar.showMessage("Arama iptal ediliyor...")
            self.search_btn.setEnabled(False)

    def add_result(self, file_path):
        self.result_list.addItem(file_path)
        
        # Dosya içeriğini okuyup satır numaralarını bul
        keywords = [k.strip() for k in self.word_edit.text().split(',') if k.strip()]
        satir_numaralari = self.get_satir_numaralari(file_path, keywords)
        
        # Sol tarafa dosyayı ekle (dosya adı + satır numaraları)
        dosya_adi = os.path.basename(file_path)
        if satir_numaralari:
            satir_str = ", ".join(map(str, satir_numaralari))
            display_text = f"{dosya_adi} (Satır: {satir_str})"
        else:
            display_text = dosya_adi
            
        list_item = QListWidgetItem(display_text)
        list_item.setData(Qt.UserRole, file_path)  # Tam yolu data olarak sakla
        self.dosya_listesi.addItem(list_item)
        
        # İlk bulunan dosyayı otomatik seç
        if self.dosya_listesi.count() == 1:
            self.dosya_listesi.setCurrentRow(0)
    
    def get_satir_numaralari(self, file_path, keywords):
        """Dosyadaki anahtar kelimelerin bulunduğu satır numaralarını döndürür"""
        try:
            ext = os.path.splitext(file_path)[1].lower()
            full_content = ""
            
            if ext == '.txt':
                with open(file_path, 'r', encoding='utf-8') as f:
                    full_content = f.read()
                    
            elif ext in ['.docx', '.docm', '.dotx', '.dotm']:
                try:
                    from docx import Document
                    doc = Document(file_path)
                    paragraphs = []
                    for p in doc.paragraphs:
                        if p.text.strip():
                            paragraphs.append(p.text.strip())
                    full_content = '\n'.join(paragraphs)
                except ImportError:
                    return []
                    
            elif ext in ['.xlsx', '.xlsm', '.xltx', '.xltm']:
                try:
                    from openpyxl import load_workbook
                    wb = load_workbook(file_path, data_only=True)
                    all_cells = []
                    for sheet in wb.worksheets:
                        for row in sheet.iter_rows(values_only=True):
                            for cell in row:
                                if cell is not None:
                                    all_cells.append(str(cell))
                    full_content = '\n'.join(all_cells)
                except ImportError:
                    return []
                    
            elif ext == '.xlsb':
                try:
                    from pyxlsb import open_workbook
                    all_cells = []
                    with open_workbook(file_path) as wb:
                        for sheet_name in wb.get_sheet_names():
                            with wb.get_sheet(sheet_name) as sheet:
                                for row in sheet.rows():
                                    for cell in row:
                                        if cell.v is not None:
                                            all_cells.append(str(cell.v))
                    full_content = '\n'.join(all_cells)
                except ImportError:
                    return []
                    
            elif ext in ['.pptx', '.pptm', '.ppsx', '.ppsm', '.potx', '.potm']:
                try:
                    from pptx import Presentation
                    prs = Presentation(file_path)
                    all_text = []
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text_frame") and shape.text_frame:
                                for paragraph in shape.text_frame.paragraphs:
                                    if paragraph.text.strip():
                                        all_text.append(paragraph.text.strip())
                    full_content = '\n'.join(all_text)
                except ImportError:
                    return []
                    
            elif ext == '.pdf':
                try:
                    with fitz.open(file_path) as doc:
                        all_text = []
                        for page in doc:
                            text = page.get_text()
                            all_text.append(text)
                        full_content = '\n'.join(all_text)
                except Exception:
                    return []
                    
            elif ext == '.vsdx':
                try:
                    metinler = read_vsdx_text(file_path)
                    full_content = '\n'.join(metinler)
                except Exception:
                    return []
                    
            else:
                return []
                
            # Satır numaralarını bul
            lines = full_content.split('\n')
            satir_numaralari = []
            case_sensitive = self.case_sensitive_cb.isChecked()
            match_type = self.match_button_group.checkedId()
            
            for line_num, line in enumerate(lines, 1):
                for keyword in keywords:
                    if self.matches_keyword(line, keyword, match_type, case_sensitive):
                        satir_numaralari.append(line_num)
                        break  # Bu satırda bir anahtar kelime bulundu, diğerlerine bakmaya gerek yok
                        
            return sorted(list(set(satir_numaralari)))  # Tekrarları kaldır ve sırala
            
        except Exception:
            return []

    def search_finished(self, count):
        if self._searching:
            if count == 0:
                # Hiç sonuç bulunamadığında pop-up uyarısı göster
                msg = QMessageBox()
                msg.setIcon(QMessageBox.Information)
                msg.setWindowTitle("Arama Sonucu")
                msg.setText("Arama tamamlandı")
                msg.setInformativeText("Aranılan kelimeler hiçbir dosyada bulunamadı.\n\nÖneriler:\n• Farklı kelimeler deneyin\n• Büyük/küçük harf ayarını kontrol edin\n• Dosya türü seçimlerini gözden geçirin")
                msg.setStandardButtons(QMessageBox.Ok)
                msg.exec_()
                self.status_bar.showMessage("Aranılan kelimeler bulunamadı.")
            else:
                self.status_bar.showMessage(f"Arama tamamlandı. {count} dosya bulundu.")
        else:
            self.status_bar.showMessage("Arama iptal edildi.")
        self.search_btn.setText("Aramayı Başlat")
        self.search_btn.setStyleSheet("background-color: #1976d2; color: white; font-weight: bold; font-size: 16px; border-radius: 8px;")
        self.search_btn.setEnabled(True)
        self._searching = False

    def open_selected_file(self, item):
        file_path = item.text()
        try:
            if platform.system() == "Windows":
                os.startfile(file_path)
            elif platform.system() == "Darwin":
                subprocess.call(["open", file_path])
            else:
                subprocess.call(["xdg-open", file_path])
        except Exception as e:
            self.status_bar.showMessage(f"Dosya açılamadı: {e}")

    def show_context_menu(self, pos):
        item = self.result_list.itemAt(pos)
        if not item:
            return
        file_path = item.text()
        menu = QMenu()
        ac_action = menu.addAction("Dosyayı Aç")
        konum_action = menu.addAction("Dosyanın Konumunu Aç")
        kopyala_action = menu.addAction("Yolu Kopyala")
        action = menu.exec_(QCursor.pos())
        if action == ac_action:
            self.open_selected_file(item)
        elif action == konum_action:
            self.open_file_location(file_path)
        elif action == kopyala_action:
            self.copy_file_path(file_path)

    def open_file_location(self, file_path):
        try:
            folder = os.path.dirname(file_path)
            if platform.system() == "Windows":
                subprocess.Popen(f'explorer /select,"{file_path}"')
            elif platform.system() == "Darwin":
                subprocess.call(["open", folder])
            else:
                subprocess.call(["xdg-open", folder])
        except Exception as e:
            self.status_bar.showMessage(f"Konum açılamadı: {e}")

    def copy_file_path(self, file_path):
        try:
            cb = QApplication.clipboard()
            cb.setText(file_path)
            self.status_bar.showMessage("Dosya yolu panoya kopyalandı.")
        except Exception as e:
            self.status_bar.showMessage(f"Kopyalama hatası: {e}")

    def save_results(self):
        if self.result_list.count() == 0:
            self.status_bar.showMessage("Kaydedilecek sonuç yok.")
            return
        path, _ = QFileDialog.getSaveFileName(self, "Sonuçları Kaydet", "", "Metin Dosyası (*.txt);;CSV Dosyası (*.csv)")
        if not path:
            return
        try:
            with open(path, 'w', encoding='utf-8') as f:
                for i in range(self.result_list.count()):
                    line = self.result_list.item(i).text()
                    if path.endswith('.csv'):
                        f.write(f'"{line}"\n')
                    else:
                        f.write(line + '\n')
            self.status_bar.showMessage(f"Sonuçlar kaydedildi: {path}")
        except Exception as e:
            self.status_bar.showMessage(f"Kayıt hatası: {e}")

    def show_preview(self, current, previous):
        # Bu fonksiyon artık dosya_secildi tarafından kullanılıyor
        pass

    def show_keyword_locations(self, content, keywords):
        lines = content.split('\n')
        result_lines = []
        case_sensitive = self.case_sensitive_cb.isChecked()
        match_type = self.match_button_group.checkedId()
        
        # Her satırı kontrol et ve anahtar kelime varsa satır numarasyla birlikte göster
        for line_num, line in enumerate(lines, 1):
            found_keywords = []
            
            # Bu satırda hangi anahtar kelimeler var
            for keyword in keywords:
                if self.matches_keyword(line, keyword, match_type, case_sensitive):
                    found_keywords.append(keyword)
            
            if found_keywords:
                # Satır numarası + bulunan kelimeler + içerik
                keywords_str = ", ".join(found_keywords)
                result_lines.append(f"[Satır {line_num}] ({keywords_str}): {line.strip()}")
        
        return '\n\n'.join(result_lines) if result_lines else ""

    def dosya_secildi(self, current_item):
        if current_item:
            # Tam dosya yolunu data'dan al
            actual_file_path = current_item.data(Qt.UserRole)
            
            if not actual_file_path:
                return
                
            self.satir_bilgileri.clear() # Satır bilgilerini temizle
            
            # Dosya var mı kontrol et
            if not os.path.exists(actual_file_path):
                return
                
            keywords = [k.strip() for k in self.word_edit.text().split(',') if k.strip()]
            
            if not keywords:
                return
                
            ext = os.path.splitext(actual_file_path)[1].lower()
            
            try:
                full_content = ""
                
                if ext == '.txt':
                    with open(actual_file_path, 'r', encoding='utf-8') as f:
                        full_content = f.read()
                        
                elif ext in ['.docx', '.docm', '.dotx', '.dotm']:
                    try:
                        from docx import Document
                        doc = Document(actual_file_path)
                        paragraphs = []
                        for p in doc.paragraphs:
                            if p.text.strip():
                                paragraphs.append(p.text.strip())
                        full_content = '\n'.join(paragraphs)
                    except ImportError:
                        return
                        
                elif ext in ['.xlsx', '.xlsm', '.xltx', '.xltm']:
                    try:
                        from openpyxl import load_workbook
                        wb = load_workbook(actual_file_path, data_only=True)
                        all_cells = []
                        for sheet in wb.worksheets:
                            for row in sheet.iter_rows(values_only=True):
                                for cell in row:
                                    if cell is not None:
                                        all_cells.append(str(cell))
                        full_content = '\n'.join(all_cells)
                    except ImportError:
                        return
                        
                elif ext == '.xlsb':
                    try:
                        from pyxlsb import open_workbook
                        all_cells = []
                        with open_workbook(actual_file_path) as wb:
                            for sheet_name in wb.get_sheet_names():
                                with wb.get_sheet(sheet_name) as sheet:
                                    for row in sheet.rows():
                                        for cell in row:
                                            if cell.v is not None:
                                                all_cells.append(str(cell.v))
                        full_content = '\n'.join(all_cells)
                    except ImportError:
                        return
                        
                elif ext in ['.pptx', '.pptm', '.ppsx', '.ppsm', '.potx', '.potm']:
                    try:
                        from pptx import Presentation
                        prs = Presentation(actual_file_path)
                        all_text = []
                        for slide in prs.slides:
                            for shape in slide.shapes:
                                if hasattr(shape, "text_frame") and shape.text_frame:
                                    for paragraph in shape.text_frame.paragraphs:
                                        if paragraph.text.strip():
                                            all_text.append(paragraph.text.strip())
                        full_content = '\n'.join(all_text)
                    except ImportError:
                        return
                        
                elif ext == '.pdf':
                    try:
                        with fitz.open(actual_file_path) as doc:
                            all_text = []
                            for page in doc:
                                text = page.get_text()
                                all_text.append(text)
                            full_content = '\n'.join(all_text)
                    except Exception:
                        return
                        
                elif ext == '.vsdx':
                    try:
                        metinler = read_vsdx_text(actual_file_path)
                        full_content = '\n'.join(metinler)
                    except Exception:
                        return
                        
                else:
                    return
                    
                # Anahtar kelime bulunan satırları bul ve sağ sütuna ekle
                lines = full_content.split('\n')
                case_sensitive = self.case_sensitive_cb.isChecked()
                match_type = self.match_button_group.checkedId()
                
                for line_num, line in enumerate(lines, 1):
                    found_keywords = []
                    
                    # Bu satırda hangi anahtar kelimeler var
                    for keyword in keywords:
                        if self.matches_keyword(line, keyword, match_type, case_sensitive):
                            found_keywords.append(keyword)
                    
                    if found_keywords:
                        # Satır numarası ve içeriği göster
                        if len(line.strip()) > 60:
                            display_text = f"Satır {line_num}: {line.strip()[:60]}..."
                        else:
                            display_text = f"Satır {line_num}: {line.strip()}"
                        self.satir_bilgileri.addItem(display_text)
                    
            except Exception:
                pass

    def matches_keyword(self, text, keyword, match_type, case_sensitive):
        """Kelime eşleştirme türüne göre arama yapar"""
        
        # Büyük/küçük harf ayarı
        search_text = text if case_sensitive else text.lower()
        search_keyword = keyword if case_sensitive else keyword.lower()
        
        if match_type == 0:  # Normal arama
            return search_keyword in search_text
        elif match_type == 1:  # Tam kelime
            # Word boundary kullanarak tam kelime ara
            flags = 0 if case_sensitive else re.IGNORECASE
            pattern = r'\b' + re.escape(keyword) + r'\b'
            return bool(re.search(pattern, text, flags))
        elif match_type == 2:  # Başlangıç
            # Kelime başında ara
            flags = 0 if case_sensitive else re.IGNORECASE
            pattern = r'\b' + re.escape(keyword)
            return bool(re.search(pattern, text, flags))
        elif match_type == 3:  # Bitiş
            # Kelime sonunda ara
            flags = 0 if case_sensitive else re.IGNORECASE
            pattern = re.escape(keyword) + r'\b'
            return bool(re.search(pattern, text, flags))
        
        return False

def read_vsdx_text(file_path):
    metinler = []
    try:
        with zipfile.ZipFile(file_path, 'r') as zf:
            for name in zf.namelist():
                if name.startswith('visio/pages/page') and name.endswith('.xml'):
                    try:
                        data = zf.read(name)
                        root = ET.parse(io.BytesIO(data)).getroot()
                        for t in root.iter('{*}t'):
                            if t.text:
                                metinler.append(t.text)
                    except Exception:
                        continue
    except zipfile.BadZipFile:
        pass
    return metinler

if __name__ == "__main__":
    app = QApplication(sys.argv)
    window = MainWindow()
    window.show()
    sys.exit(app.exec_()) 
